# theme_analyzer.py
# -*- coding: utf-8 -*-

"""
Part 1: Presentation Theme and Style Analyzer

This script defines the ThemeAnalyzer class, which is responsible for extracting
the design and style elements from a given PowerPoint (.pptx) file. It identifies
common fonts, sizes, and colors for titles and body text, as well as the primary
and accent colors from the presentation's theme.

Dependencies:
- python-pptx: To interact with PowerPoint files.
- collections.Counter: For finding the most common elements.

To install dependencies:
pip install python-pptx
"""

import collections
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import PP_PLACEHOLDER

class ThemeAnalyzer:
    """
    Analyzes a PowerPoint presentation to extract its theme and style information.
    """
    def __init__(self, pptx_path):
        """
        Initializes the ThemeAnalyzer with the path to a .pptx file.

        Args:
            pptx_path (str): The file path to the PowerPoint presentation.
        """
        try:
            self.prs = Presentation(pptx_path)
        except Exception as e:
            print(f"Error opening presentation file: {e}")
            self.prs = None
            return

        self.title_font = None
        self.body_font = None
        self.primary_color = None
        self.accent_color = None
        self._analyze_theme()

    def _analyze_theme(self):
        """
        Performs the analysis of the presentation's theme and styles.
        """
        if not self.prs or not self.prs.slide_masters:
            print("Presentation contains no slide masters, cannot analyze theme.")
            return

        title_fonts, body_fonts = [], []
        title_font_sizes, body_font_sizes = [], []
        title_colors, body_colors = [], []
        
        # Analyze slide masters for common fonts, sizes, and colors by checking placeholder types.
        for slide_master in self.prs.slide_masters:
            for placeholder in slide_master.placeholders:
                # Check for various types of title placeholders
                if placeholder.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE]:
                    if hasattr(placeholder, 'text_frame') and placeholder.text_frame:
                        font = placeholder.text_frame.paragraphs[0].font
                        if font.name:
                            title_fonts.append(font.name)
                        if font.size:
                            title_font_sizes.append(font.size.pt)
                        # --- FIX: Check for .rgb attribute before accessing it ---
                        if hasattr(font.color, 'rgb') and font.color.rgb is not None:
                            title_colors.append(font.color.rgb)
                
                # Check for Body placeholders
                elif placeholder.placeholder_format.type == PP_PLACEHOLDER.BODY:
                     if hasattr(placeholder, 'text_frame') and placeholder.text_frame:
                        font = placeholder.text_frame.paragraphs[0].font
                        if font.name:
                            body_fonts.append(font.name)
                        if font.size:
                            body_font_sizes.append(font.size.pt)
                        # --- FIX: Check for .rgb attribute before accessing it ---
                        if hasattr(font.color, 'rgb') and font.color.rgb is not None:
                            body_colors.append(font.color.rgb)

        # --- CORRECTED & ROBUST FONT/SIZE/COLOR DETERMINATION ---
        title_font_counter = collections.Counter(title_fonts)
        title_size_counter = collections.Counter(title_font_sizes)
        body_font_counter = collections.Counter(body_fonts)
        body_size_counter = collections.Counter(body_font_sizes)
        title_color_counter = collections.Counter(title_colors)
        body_color_counter = collections.Counter(body_colors)

        self.title_font = {
            'name': title_font_counter.most_common(1)[0][0] if title_font_counter else 'Calibri',
            'size': title_size_counter.most_common(1)[0][0] if title_size_counter else 32
        }
        self.body_font = {
            'name': body_font_counter.most_common(1)[0][0] if body_font_counter else 'Calibri',
            'size': body_size_counter.most_common(1)[0][0] if body_size_counter else 18
        }
        
        # The primary color is the most common title color.
        self.primary_color = title_color_counter.most_common(1)[0][0] if title_color_counter else RGBColor(0, 0, 0)
        # The accent color is the most common body text color.
        self.accent_color = body_color_counter.most_common(1)[0][0] if body_color_counter else RGBColor(89, 89, 89)


    def get_style(self):
        """
        Returns a dictionary of the detected theme properties.
        """
        if not self.prs:
            return None
            
        return {
            "title_font_name": self.title_font['name'],
            "title_font_size": self.title_font['size'],
            "body_font_name": self.body_font['name'],
            "body_font_size": self.body_font['size'],
            "primary_color_rgb": str(self.primary_color),
            "accent_color_rgb": str(self.accent_color)
        }

if __name__ == '__main__':
    # This is an example of how to use the ThemeAnalyzer class.
    # To run this, you need a PowerPoint file named 'sample.pptx' in the same directory.
    # You can create a simple presentation and save it with this name.
    
    try:
        # Create a dummy presentation for testing if one doesn't exist.
        from pptx.util import Inches
        
        try:
            prs = Presentation('sample.pptx')
        except:
            prs = Presentation()
            # Use a layout that has a title and body
            title_and_content_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(title_and_content_layout)
            
            # Set text to ensure placeholders are populated for analysis
            title = slide.shapes.title
            title.text = "Test Title"
            
            body = slide.placeholders[1]
            body.text = "This is a test bullet point."
            
            prs.save('sample.pptx')
            print("Created a dummy 'sample.pptx' for testing.")

        analyzer = ThemeAnalyzer('sample.pptx')
        style_guide = analyzer.get_style()

        if style_guide:
            print("--- Detected Presentation Style ---")
            print(f"Title Font: {style_guide['title_font_name']}, Size: {style_guide['title_font_size']}pt")
            print(f"Body Font: {style_guide['body_font_name']}, Size: {style_guide['body_font_size']}pt")
            print(f"Primary Color (RGB): {style_guide['primary_color_rgb']}")
            print(f"Accent Color (RGB): {style_guide['accent_color_rgb']}")
            print("---------------------------------")
        else:
            print("Could not analyze the presentation.")

    except ImportError:
        print("Please install python-pptx: pip install python-pptx")
    except Exception as e:
        print(f"An error occurred: {e}")
