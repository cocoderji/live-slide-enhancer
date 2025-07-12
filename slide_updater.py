# slide_updater.py
# -*- coding: utf-8 -*-

"""
Part 4: Live Slide Update Mechanism (Data-First Logic)

This version prioritizes generating a data chart over fetching an image
if the content is data-driven.
"""

from pptx import Presentation
import win32com.client
import os

# Constants
msoTrue = -1
ppLayoutBlank = 12
ppAlignCenter = 2
ppAlignLeft = 1

class SlideUpdater:
    def __init__(self, pptx_path):
        self.pptx_path = os.path.abspath(pptx_path)

    def _get_active_presentation(self):
        try:
            app = win32com.client.GetActiveObject("PowerPoint.Application")
            for pres in app.Presentations:
                if pres.FullName.lower() == self.pptx_path.lower():
                    if pres.SlideShowWindow: return pres
            return None
        except Exception: return None

    def start_presentation_show(self):
        try:
            app = win32com.client.Dispatch("PowerPoint.Application")
            app.Visible = msoTrue
            presentation = app.Presentations.Open(self.pptx_path)
            presentation.SlideShowSettings.Run()
            return True
        except Exception as e:
            print(f"ERROR: Failed to open or start the slideshow: {e}"); return False

    def get_current_slide_index(self):
        active_presentation = self._get_active_presentation()
        if active_presentation:
            try: return active_presentation.SlideShowWindow.View.CurrentShowPosition
            except Exception: return None
        return None

    def get_text_from_slide(self, slide_index):
        if not os.path.exists(self.pptx_path): return ""
        all_text = []
        try:
            prs = Presentation(self.pptx_path)
            if not (0 < slide_index <= len(prs.slides)): return ""
            slide = prs.slides[slide_index - 1]
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text:
                    all_text.append(shape.text_frame.text)
            return "\n".join(all_text)
        except Exception as e:
            print(f"ERROR: Could not extract text from slide {slide_index}: {e}"); return ""

    def _add_content_to_slide(self, target_slide, content_json, style_guide, visual_generator):
        """A helper function to add styled text and visuals to a given slide."""
        layout = content_json.get('layout', 'text_left_visual_right')
        print(f"INFO: Building slide with layout: '{layout}'")

        # --- NEW: Prioritize chart data over other visuals ---
        visual_path = None
        if content_json.get('chart_data'):
            visual_path = visual_generator.create_chart(content_json['chart_data'])
        elif content_json.get('image_suggestion'):
            visual_path = visual_generator.get_icon(content_json['image_suggestion'])
            if not visual_path:
                visual_path = visual_generator.get_image(content_json['image_suggestion'])
        
        # Determine text box width based on whether a visual will be present
        inch_to_points = lambda i: i * 72
        text_box_width = inch_to_points(5.5) if visual_path else inch_to_points(9)

        # Add Title
        title_box = target_slide.Shapes.AddTextbox(1, inch_to_points(0.5), inch_to_points(0.2), text_box_width, inch_to_points(1.5))
        title_range = title_box.TextFrame.TextRange
        title_range.Text = content_json.get('title', 'New Topic')
        title_font = title_range.Font
        title_font.Name = style_guide.get('title_font_name', 'Calibri')
        title_font.Size = int(style_guide.get('title_font_size', 32))
        title_font.Bold = msoTrue
        title_font.Color.RGB = int(style_guide['primary_color_rgb'], 16)
        title_range.ParagraphFormat.Alignment = ppAlignCenter if not visual_path else ppAlignLeft

        # Add Body
        body_box = target_slide.Shapes.AddTextbox(1, inch_to_points(0.5), inch_to_points(1.8), text_box_width, inch_to_points(5))
        body_frame = body_box.TextFrame
        body_frame.TextRange.Text = ""
        for point_text in content_json.get('points', []):
            p = body_frame.TextRange.InsertAfter(point_text + "\r\n")
            p.ParagraphFormat.Bullet.Visible = msoTrue
            p.Font.Name = style_guide.get('body_font_name', 'Calibri')
            p.Font.Size = int(style_guide.get('body_font_size', 18))
            p.Font.Color.RGB = int(style_guide['accent_color_rgb'], 16)
        body_frame.TextRange.ParagraphFormat.Alignment = ppAlignLeft
        
        # Add Branding Watermark
        logo_box = target_slide.Shapes.AddTextbox(1, inch_to_points(0.2), inch_to_points(7.2), inch_to_points(2), inch_to_points(0.5))
        logo_box.TextFrame.TextRange.Text = "Updated live by Savi.ai"
        logo_font = logo_box.TextFrame.TextRange.Font
        logo_font.Size = 10
        logo_font.Color.RGB = int("A0A0A0", 16)
        
        # Add the visual if one was created
        if visual_path and os.path.exists(visual_path):
            target_slide.Shapes.AddPicture(FileName=os.path.abspath(visual_path), LinkToFile=False, SaveWithDocument=True, 
                                           Left=inch_to_points(6), Top=inch_to_points(2.5), 
                                           Width=inch_to_points(3.5), Height=inch_to_points(3.5))

    def insert_new_slide_after_current(self, slide_index, content_json, style_guide, visual_generator):
        active_presentation = self._get_active_presentation()
        if not active_presentation: print("ERROR: Could not find active presentation."); return
        try:
            new_slide_index = slide_index + 1
            target_slide = active_presentation.Slides.Add(new_slide_index, ppLayoutBlank)
            self._add_content_to_slide(target_slide, content_json, style_guide, visual_generator)
            active_presentation.SlideShowWindow.View.GotoSlide(new_slide_index)
            active_presentation.Save()
        except Exception as e:
            print(f"ERROR: An error occurred during slide insertion: {e}")

    def update_existing_slide(self, slide_index, content_json, style_guide, visual_generator):
        active_presentation = self._get_active_presentation()
        if not active_presentation: print("ERROR: Could not find active presentation."); return
        try:
            target_slide = active_presentation.Slides(slide_index)
            for i in range(target_slide.Shapes.Count, 0, -1): target_slide.Shapes(i).Delete()
            self._add_content_to_slide(target_slide, content_json, style_guide, visual_generator)
            active_presentation.SlideShowWindow.View.GotoSlide(slide_index)
            active_presentation.Save()
        except Exception as e:
            print(f"ERROR: An error occurred during slide update: {e}")
